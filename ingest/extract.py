#!/usr/bin/env python3
"""
ingest/extract.py

Production-minded document extractor for RAG ingestion pipeline.

Features:
- Processes files from an input directory (or single file).
- Supports PDF, DOCX, PPTX, TXT, CSV/XLSX, and common image formats.
- Extracts page-level text (PDF/PPTX) or row-level text (CSV/XLSX).
- Extracts images from PDFs/pages and runs OCR on images.
- Produces structured JSON per input file in output directory.
- Idempotent: uses file-hash + output file existence to avoid reprocessing.
- Configurable via CLI args (input_dir, output_dir, min_text_length, upload_to_s3).
- Hooks: easy to add post-processing (send message to Kafka / enqueue embedding job).
"""

import argparse
import hashlib
import json
import logging
import os
from pathlib import Path
from typing import Dict, List, Optional, Any
from PIL import Image, ImageEnhance

# Optional dependencies; import only when needed to fail gracefully
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import pytesseract
    from PIL import Image, ImageEnhance
except Exception:
    pytesseract = None
    Image = None
    ImageEnhance = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pandas as pd
except Exception:
    pd = None

# Basic logging setup
logging.basicConfig(
    level=os.environ.get("EXTRACT_LOG_LEVEL", "INFO"),
    format="%(asctime)s %(levelname)s %(name)s - %(message)s",
)
logger = logging.getLogger("extractor")


# -------------------------
# Utility helpers
# -------------------------
def compute_sha256(path: Path, block_size: int = 65536) -> str:
    """Compute SHA256 of a file (streaming)."""
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(block_size), b""):
            h.update(block)
    return h.hexdigest()


def safe_write_json(obj: Any, out_path: Path):
    """Write JSON atomically (write to tmp then rename)."""
    tmp = out_path.with_suffix(out_path.suffix + ".tmp")
    with tmp.open("w", encoding="utf-8") as f:
        json.dump(obj, f, ensure_ascii=False, indent=2)
    tmp.replace(out_path)


def ocr_image_pil(img: Image.Image) -> str:
    """Run lightweight preprocessing + pytesseract OCR on PIL Image."""
    if pytesseract is None or Image is None:
        raise RuntimeError("pytesseract and Pillow required for OCR")
    try:
        # grayscale + increase contrast + binary threshold - simple but effective
        img_gray = img.convert("L")
        img_enh = ImageEnhance.Contrast(img_gray).enhance(1.8)
        img_bw = img_enh.point(lambda x: 0 if x < 128 else 255, "1")
        text = pytesseract.image_to_string(img_bw, lang="eng", config="--psm 6")
        return text.strip()
    except Exception as e:
        logger.exception("OCR failed: %s", e)
        return ""


# -------------------------
# Extractors by file type
# -------------------------
def extract_pdf(pdf_path: Path, min_text_length: int = 20) -> Optional[Dict]:
    """Extract text and images from a PDF using pdfplumber. Returns structured dict."""
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is required to extract PDFs. Install pdfplumber.")
    logger.info("Processing PDF: %s", pdf_path)
    out = {"file_name": pdf_path.name, "file_hash": compute_sha256(pdf_path), "type": "pdf", "pages": []}

    try:
        with pdfplumber.open(str(pdf_path)) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages, start=1):
                try:
                    text = page.extract_text() or ""
                    # Save a rasterized page image for OCR or image cropping
                    # pdfplumber page.to_image() can create a PIL image
                    page_img = None
                    try:
                        page_img = page.to_image(resolution=150).original
                    except Exception:
                        page_img = None

                    ocr_text = ""
                    image_path = None
                    if page_img is not None and pytesseract is not None:
                        ocr_text = ocr_image_pil(page_img)

                        # Optionally save page image for debugging / multimodal pipeline
                        images_dir = pdf_path.parent / "extracted_images"
                        images_dir.mkdir(parents=True, exist_ok=True)
                        image_path = images_dir / f"{pdf_path.stem}_page{i}.png"
                        page_img.save(image_path)

                    # Skip pages that are essentially empty (both text and OCR below threshold)
                    if len(text.strip()) < min_text_length and len(ocr_text.strip()) < min_text_length:
                        logger.debug("Skipping mostly-empty page %d/%d for %s", i, total_pages, pdf_path.name)
                        continue

                    page_record = {
                        "page_number": i,
                        "text": text.strip(),
                        "ocr_text": ocr_text,
                        "image_path": str(image_path) if image_path else None,
                    }
                    out["pages"].append(page_record)
                except Exception as e:
                    logger.exception("Failed to process page %d of %s: %s", i, pdf_path, e)
                    continue
    except Exception as e:
        logger.exception("Failed to open PDF %s: %s", pdf_path, e)
        return None

    if not out["pages"]:
        logger.warning("No meaningful pages extracted for %s", pdf_path)
        return None
    return out


def extract_docx(docx_path: Path, min_text_length: int = 20) -> Optional[Dict]:
    """Extract text from a DOCX file (python-docx)."""
    if docx is None:
        raise RuntimeError("python-docx is required to extract .docx")
    logger.info("Processing DOCX: %s", docx_path)
    doc = docx.Document(str(docx_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and len(p.text.strip()) >= min_text_length]
    if not paragraphs:
        logger.warning("No text found in DOCX: %s", docx_path)
        return None
    return {"file_name": docx_path.name, "file_hash": compute_sha256(docx_path), "type": "docx", "paragraphs": paragraphs}


def extract_pptx(pptx_path: Path, min_text_length: int = 20) -> Optional[Dict]:
    """Extract text from PPTX slides (python-pptx)."""
    if Presentation is None:
        raise RuntimeError("python-pptx is required to extract .pptx")
    logger.info("Processing PPTX: %s", pptx_path)
    prs = Presentation(str(pptx_path))
    slides_out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t and len(t) >= min_text_length:
                    texts.append(t)
        if texts:
            slides_out.append({"slide_number": i, "texts": texts})
    if not slides_out:
        logger.warning("No text found in PPTX: %s", pptx_path)
        return None
    return {"file_name": pptx_path.name, "file_hash": compute_sha256(pptx_path), "type": "pptx", "slides": slides_out}


def extract_text_file(txt_path: Path, min_text_length: int = 20) -> Optional[Dict]:
    """Extract plain text file."""
    logger.info("Processing TXT: %s", txt_path)
    s = txt_path.read_text(encoding="utf-8", errors="ignore").strip()
    if len(s) < min_text_length:
        logger.warning("Text file too short: %s", txt_path)
        return None
    return {"file_name": txt_path.name, "file_hash": compute_sha256(txt_path), "type": "txt", "text": s}


def extract_csv_xlsx(path: Path, min_text_length: int = 1) -> Optional[Dict]:
    """Extract CSV or Excel - returns rows converted to strings for indexing."""
    if pd is None:
        raise RuntimeError("pandas is required for CSV/XLSX extraction")
    logger.info("Processing tabular file: %s", path)
    try:
        if path.suffix.lower() in [".csv"]:
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
    except Exception as e:
        logger.exception("Failed to read tabular file %s: %s", path, e)
        return None

    # Convert rows to textual representation; you'll likely want a better table extractor later
    rows = []
    for row in df.itertuples(index=False, name=None):
        # convert tuple row into a joined string
        row_str = " | ".join("" if v is None else str(v) for v in row)
        if len(row_str.strip()) >= min_text_length:
            rows.append(row_str)
    if not rows:
        logger.warning("No rows (after filtering) for %s", path)
        return None
    return {"file_name": path.name, "file_hash": compute_sha256(path), "type": "table", "rows": rows}


# -------------------------
# Orchestration / file loop
# -------------------------
SUPPORTED = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_text_file,
    ".csv": extract_csv_xlsx,
    ".xlsx": extract_csv_xlsx,
    ".xls": extract_csv_xlsx,
    # images - treat single images as trivially OCR-able "pages"
    ".png": lambda p, **kw: {"file_name": p.name, "file_hash": compute_sha256(p), "type": "image", "ocr_text": ocr_image_pil(Image.open(p))} if Image else None,
    ".jpg": lambda p, **kw: {"file_name": p.name, "file_hash": compute_sha256(p), "type": "image", "ocr_text": ocr_image_pil(Image.open(p))} if Image else None,
    ".jpeg": lambda p, **kw: {"file_name": p.name, "file_hash": compute_sha256(p), "type": "image", "ocr_text": ocr_image_pil(Image.open(p))} if Image else None,
}


def process_path(input_path: Path, output_dir: Path, min_text_length: int = 20, skip_if_exists: bool = True) -> Optional[Path]:
    """Process a single input file and write JSON to output_dir. Returns output path or None."""
    logger.info("Starting processing for %s", input_path)
    suffix = input_path.suffix.lower()
    extractor = SUPPORTED.get(suffix)
    if extractor is None:
        logger.warning("Unsupported file type: %s (skipping)", input_path)
        return None

    # idempotency: check existing output by hash
    file_hash = compute_sha256(input_path)
    out_name = f"{input_path.stem}_{file_hash[:10]}.extracted.json"
    out_path = output_dir / out_name
    if out_path.exists() and skip_if_exists:
        logger.info("Output already exists for %s -> %s (skipping)", input_path.name, out_path.name)
        return out_path

    # run extractor
    try:
        extracted = extractor(input_path, min_text_length=min_text_length)
    except TypeError:
        # fallback: extractor might not accept min_text_length (image lambdas)
        extracted = extractor(input_path)

    if not extracted:
        logger.warning("Extraction returned no content for %s", input_path)
        return None

    # enrich with provenance
    extracted_meta = {
        "source_path": str(input_path),
        "file_name": input_path.name,
        "file_hash": file_hash,
        "extracted_at": __import__("datetime").datetime.utcnow().isoformat() + "Z",
        "content": extracted,
    }

    # write JSON atomically
    safe_write_json(extracted_meta, out_path)
    logger.info("Wrote extraction output to %s", out_path)
    return out_path


def process_directory(input_dir: Path, output_dir: Path, min_text_length: int = 20, skip_if_exists: bool = True) -> List[Path]:
    """Process all files under input_dir (non-recursive). Returns list of written output paths."""
    written = []
    for p in sorted(input_dir.iterdir()):
        if p.is_file():
            try:
                out = process_path(p, output_dir, min_text_length=min_text_length, skip_if_exists=skip_if_exists)
                if out:
                    written.append(out)
            except Exception as e:
                logger.exception("Failed to process %s: %s", p, e)
    return written


# -------------------------
# CLI
# -------------------------
def build_arg_parser():
    p = argparse.ArgumentParser(description="Document extraction for RAG ingestion (PDF, DOCX, PPTX, CSV, XLSX, images).")
    p.add_argument("--input", "-i", type=str, default="data", help="Input directory or file to process")
    p.add_argument("--output", "-o", type=str, default="outputs", help="Output directory for JSON files")
    p.add_argument("--min-text-length", type=int, default=20, help="Minimum text length to consider a page/paragraph valid")
    p.add_argument("--skip-if-exists", action="store_true", help="Skip files that already have output (idempotency)")
    p.add_argument("--single", action="store_true", help="Input is a single file (not a directory)")
    return p


def main():
    args = build_arg_parser().parse_args()
    input_path = Path(args.input)
    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    if args.single:
        if not input_path.exists():
            logger.error("Input file does not exist: %s", input_path)
            return
        process_path(input_path, output_dir, min_text_length=args.min_text_length, skip_if_exists=args.skip_if_exists)
    else:
        if not input_path.exists():
            logger.error("Input directory does not exist: %s", input_path)
            return
        written = process_directory(input_path, output_dir, min_text_length=args.min_text_length, skip_if_exists=args.skip_if_exists)
        logger.info("Finished processing directory. Wrote %d files.", len(written))


if __name__ == "__main__":
    main()
